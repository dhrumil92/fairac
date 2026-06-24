from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Image Paths
img_problem = r"C:\Users\dhrum\.gemini\antigravity-ide\brain\4ee4dcfb-4003-4498-b695-643142f1936c\ppt_problem_1782278416362.png"
img_solution = r"C:\Users\dhrum\.gemini\antigravity-ide\brain\4ee4dcfb-4003-4498-b695-643142f1936c\ppt_solution_1782278428543.png"
img_dashboard = r"C:\Users\dhrum\.gemini\antigravity-ide\brain\4ee4dcfb-4003-4498-b695-643142f1936c\ppt_dashboard_1782278439329.png"
img_hardware = r"C:\Users\dhrum\.gemini\antigravity-ide\brain\4ee4dcfb-4003-4498-b695-643142f1936c\ppt_hardware_1782278453734.png"

# Helper function to add a slide with text on left and image on right
def add_split_slide(prs, title_text, bullet_points, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only layout
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Text Box (Left Side)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.space_before = Pt(12)
            
    # Image (Right Side)
    try:
        slide.shapes.add_picture(img_path, Inches(5.2), Inches(1.8), width=Inches(4.3))
    except Exception as e:
        print(f"Error loading image {img_path}: {e}")

# --- Slide 1: Title Slide ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "FairAC: Smart & Fair AC Billing for Hostels"
subtitle.text = "IoT-Powered Transparent Billing System\nPitch to Hostel Management"

# --- Slide 2: The Problem ---
add_split_slide(prs, "The Problem: Unfair Billing & Wastage", [
    "• Flat Fees: Fixed monthly AC fees are unfair to students who use it less.",
    "• Roommate Disputes: Constant arguments over who left the AC on and who pays.",
    "• Massive Wastage: Students leave ACs running while in class because they 'already paid'.",
    "• No Analytics: Management has zero visibility into per-room energy consumption."
], img_problem)

# --- Slide 3: The Solution ---
add_split_slide(prs, "The Solution: FairAC System", [
    "• IoT Smart Meter: A low-cost Wi-Fi device installed on the AC switchboard.",
    "• Session Tracking: Students use the App to 'Start a Session' and add present roommates.",
    "• Automated Splitting: The exact kWh consumed is split evenly ONLY among the present roommates.",
    "• Digital Wallet: The AC shuts off automatically if their prepaid wallet runs out."
], img_solution)

# --- Slide 4: Benefits for Management ---
add_split_slide(prs, "Benefits for Hostel Management", [
    "• Zero Manual Work: No more reading sub-meters at the end of every month.",
    "• Advance Payment: Wallet system collects money BEFORE electricity is used.",
    "• Energy Conservation: Students turn off the AC to save their own money.",
    "• Admin Dashboard: Live visibility into active sessions and monthly revenue."
], img_dashboard)

# --- Slide 5: Hardware Costing ---
add_split_slide(prs, "Hardware Costing (Per Room)", [
    "• ESP32 Microcontroller: ~₹350",
    "• ACS712 Current Sensor: ~₹150",
    "• ZMPT101B Voltage Sensor: ~₹100",
    "• 30A Relay Module: ~₹200",
    "• OLED Display (0.96\"): ~₹150",
    "• Enclosure & Wiring: ~₹100",
    "",
    "Total Estimated Cost: ~₹1,050 per room"
], img_hardware)

# --- Slide 6: Pilot Deployment ---
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Next Steps: Pilot Deployment Proposal"
tf = body_shape.text_frame
tf.text = "We propose a low-risk pilot program to prove the value:"
p = tf.add_paragraph()
p.text = "Phase 1: Install the FairAC hardware in 5 rooms (Cost: ~₹5,250)."
p.level = 1
p = tf.add_paragraph()
p.text = "Phase 2: Onboard the students of those rooms onto the platform."
p.level = 1
p = tf.add_paragraph()
p.text = "Phase 3: Run the system for 1 month and compare energy usage and revenue collection with traditional rooms."
p.level = 1
p = tf.add_paragraph()
p.text = "Phase 4: Full roll-out across the entire hostel."
p.level = 1

prs.save('FairAC_Pitch_V2.pptx')
print("Successfully generated visually enhanced FairAC_Pitch_V2.pptx")
